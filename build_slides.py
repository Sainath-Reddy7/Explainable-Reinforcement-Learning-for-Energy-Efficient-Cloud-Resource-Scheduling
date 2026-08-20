"""
build_slides.py — Implementation deck for the project defense.
-> Implementation_Slides.pptx  (16:9, project palette, embedded real figures)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
R = HERE / "results"
OUT = HERE / "Implementation_Slides.pptx"

NAVY  = RGBColor(0x0F, 0x28, 0x47)
BLUE  = RGBColor(0x00, 0x66, 0xCC)
TEAL  = RGBColor(0x17, 0xA2, 0xB8)
GREEN = RGBColor(0x2D, 0x9B, 0x6B)
ORANGE= RGBColor(0xE2, 0x7D, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x77, 0x85)
LIGHT = RGBColor(0xEB, 0xF3, 0xFA)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, leading=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = runs if isinstance(runs, list) else [runs]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        if isinstance(line, tuple):          # (text, overrides dict)
            t, ov = line
        else:
            t, ov = line, {}
        r = p.add_run(); r.text = t
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
        r.font.name = "Calibri"
    return tb


def header(s, title, sub=None):
    bar = rect(s, 0, 0, prs.slide_width, Inches(1.0), NAVY)
    text(s, Inches(0.5), Inches(0.12), Inches(12.4), Inches(0.5),
         title, size=26, color=WHITE, bold=True)
    if sub:
        text(s, Inches(0.5), Inches(0.62), Inches(12.4), Inches(0.35),
             sub, size=12, color=RGBColor(0xC9, 0xD6, 0xE8))
    return s


def chip(s, x, y, w, h, title, body, accent=BLUE):
    rect(s, x, y, w, h, LIGHT, accent)
    text(s, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.35),
         title, size=13, color=accent, bold=True)
    text(s, x + Inches(0.15), y + Inches(0.42), w - Inches(0.3), h - Inches(0.5),
         body, size=10.5, color=NAVY, leading=1.05)


def bullets(s, x, y, w, h, items, size=13, gap=1.15):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = gap
        head, _, rest = it.partition("|")
        r = p.add_run(); r.text = "▪ " + head.strip() + "  "
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = NAVY; r.font.name = "Calibri"
        if rest:
            r2 = p.add_run(); r2.text = rest.strip()
            r2.font.size = Pt(size - 1); r2.font.color.rgb = MUTED
            r2.font.name = "Calibri"
    return tb


# ================================================================ #
# 1 · TITLE
# ================================================================ #
s = slide()
rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
rect(s, Inches(0.8), Inches(2.0), Inches(1.2), Inches(0.06), TEAL)
text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.6),
     [("RL-MOTS-XAI v2", {"size": 44, "bold": True, "color": WHITE}),
      ("Explainable Reinforcement Learning for Energy-Efficient Cloud Resource Scheduling",
       {"size": 20, "color": RGBColor(0x9F, 0xC5, 0xE8)})], leading=1.1)
text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.5),
     "Dueling Double-DQN + Safety Action Mask  ·  4 XAI Methods  ·  6 Audited Trust Metrics",
     size=15, color=TEAL)
text(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.9),
     ["Trained on real Google Borg traces (328 MB · 405k events) and validated on KTH SP2 HPC traces",
      "100% pure NumPy — every backprop step, Shapley value and trust metric hand-built"],
     size=13, color=RGBColor(0xC9, 0xD6, 0xE8), leading=1.3)
text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
     "Implementation Review — Integrated Cloud Computing", size=12,
     color=RGBColor(0x8F, 0xA6, 0xC0))

# ================================================================ #
# 2 · IMPLEMENTATION FLOW (6 phases)
# ================================================================ #
s = header(slide(), "Implementation Flow — Six-Phase Pipeline",
           "Every phase is a standalone module; the whole pipeline reproduces with one command")
phases = [
    ("1 · Data Ingestion", "borg_loader.py / swf_loader.py",
     "Parse 328 MB Google Borg CSV (405k events) + KTH SP2 SWF log; filter completed jobs; cache 5,000-task pools (.npz) for <0.1 s reloads", BLUE),
    ("2 · Environment", "env.py",
     "8 heterogeneous VMs (2 edge + 6 cloud) with parallel time-shared execution; 45-dim state (5 task + 5×8 VM features); QoS-dominated reward with potential-based shaping", TEAL),
    ("3 · Decision Engine", "qnetwork.py / dqn_agent.py",
     "Dueling Double-DQN in pure NumPy: shared backbone 45→128→64, value stream →32→1, advantage stream →32→8; safety mask blocks VMs >90% util before argmax", GREEN),
    ("4 · Training", "train.py",
     "300 episodes · curriculum batching 150–500 tasks · SumTree PER (50k, α=0.6, β→1) · soft target sync τ=0.005 · Adam with grad clipping", ORANGE),
    ("5 · Explainability", "explainability.py",
     "KernelSHAP (250 coalitions + WLS), Gradient×Input (exact analytic), Occlusion, Integrated Gradients — all explain the same masked policy that acts", BLUE),
    ("6 · Trust Audit + UI", "fidelity.py / dashboard",
     "6 metrics (del/ins AOPC, top-k fidelity, consistency, infidelity, stability) + operator console, publication Figure 1, animated walkthrough", TEAL),
]
for i, (t, f, b, c) in enumerate(phases):
    x = Inches(0.45 + (i % 3) * 4.25)
    y = Inches(1.35 + (i // 3) * 2.95)
    chip(s, x, y, Inches(4.05), Inches(2.75), f"{t}   ({f})", b, c)

# ================================================================ #
# 3 · MODEL ARCHITECTURE (Figure 1)
# ================================================================ #
s = header(slide(), "Model Architecture — Figure 1 (publication format)",
           "Nature-style figure: 183 mm double-column, vector PDF/SVG, zero collisions machine-verified")
s.shapes.add_picture(str(R / "fig1_architecture.png"), Inches(0.55), Inches(1.25),
                     width=Inches(12.2))
text(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.4),
     "a network · b masked dispatch · c prioritized-replay training · d audited attributions (real deletion-AOPC values)",
     size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ================================================================ #
# 4 · KEY IMPLEMENTATION DETAILS
# ================================================================ #
s = header(slide(), "Key Implementation Details", "What makes this build work")
bullets(s, Inches(0.55), Inches(1.35), Inches(6.3), Inches(5.9), [
    "Dueling decomposition | Q(s,a) = V(s) + (A − Ā): value stream learns pool health, advantage stream learns per-VM preference — ablation: removing it costs +34% cost",
    "Safety action mask | pure function of the observed state: util > 0.9 ⇒ Q ← −∞ before argmax AND during exploration — ablation: +37% cost without it",
    "Double DQN + PER | online net selects, target net evaluates (kills overestimation); SumTree samples high-|TD| transitions 50k× faster than a list scan",
    "Potential-based shaping | F = γΦ(s′) − Φ(s) with Φ = −(imbalance + queue): dense signal, provably policy-invariant (Ng et al. 1999)",
], size=13)
bullets(s, Inches(7.1), Inches(1.35), Inches(5.7), Inches(5.9), [
    "Parallel VM model | each VM tracks its live task set; finish time stretched by oversubscription ratio → schedulers produce genuinely distinct makespans/miss rates",
    "Corrected AOPC | deletion/insertion normalized by Q-spread (decision margin) instead of |Q| — fixes sign-flip on large Q, clamped to [−5, 5]",
    "Explain the deployed policy | XAI scores Q + mask — the exact function that acts — so fidelity measures what operators actually see",
    "One-command reproducibility | config.yaml drives everything; mtime-keyed caches; manifest logs seed, platform, dataset stats, runtime",
], size=13)

# ================================================================ #
# 5 · EVALUATION SETUP + RESULTS
# ================================================================ #
s = header(slide(), "Evaluation & Results — Real Borg Workloads",
           "8 schedulers × 5 loads (200–1000 tasks) × 5 seeds; identical VM pool (seed 42) for fairness")
rows = [
    ("Scheduler", "Makespan (s)", "Cost ($)", "Miss", "DI"),
    ("Min-Min (best heuristic)", "1194 ± 18", "32,036", "0.0%", "4.29"),
    ("DQN (ours)", "1194 ± 18", "29,141", "0.0%", "1.04"),
]
tbl = s.shapes.add_table(3, 5, Inches(0.55), Inches(1.4), Inches(6.4), Inches(1.5)).table
for j, wdt in enumerate([2.3, 1.3, 1.1, 0.85, 0.85]):
    tbl.columns[j].width = Inches(wdt)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        for r in para.runs:
            r.font.size = Pt(11); r.font.name = "Calibri"
            r.font.bold = i <= 1 or (i == 2 and True)
            r.font.color.rgb = WHITE if i == 0 else (GREEN if i == 2 else NAVY)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if i == 0 else (RGBColor(0xE8, 0xF7, 0xEE) if i == 2 else LIGHT)
bullets(s, Inches(0.55), Inches(3.15), Inches(6.4), Inches(3.9), [
    "Ties the best makespan & miss rate across every load",
    "Significantly cheaper than Min-Min (Wilcoxon p = 0.0004, n = 75)",
    "#1 load balance of all 8 schedulers (DI 1.04 ± 0.03)",
    "Cross-dataset: Borg's champion Min-Min collapses 14× on KTH; our DQN stays within 1.44× of best",
], size=12.5)
s.shapes.add_picture(str(R / "scheduler_comparison.png"), Inches(7.25), Inches(1.35),
                     height=Inches(5.4))

# ================================================================ #
# 6 · XAI + TRUST (the contribution)
# ================================================================ #
s = header(slide(), "Explainability Layer & Trust Audit — The Contribution",
           "Explanations are mathematically audited, not just displayed")
bullets(s, Inches(0.55), Inches(1.35), Inches(6.2), Inches(3.1), [
    "4 attribution methods | KernelSHAP · Gradient×Input · Occlusion · Integrated Gradients explain the same 60 live dispatch decisions",
    "6 trust metrics | deletion/insertion AOPC · top-k fidelity · consistency · infidelity · stability — every method verified faithful (positive AOPC)",
    "Audited trade-off | Occlusion most faithful (1.13) vs Gradient×Input 340× faster (0.16 ms) — no single method dominates",
], size=12.5)
s.shapes.add_picture(str(R / "xai_method_comparison.png"), Inches(0.55), Inches(4.35),
                     width=Inches(6.2))
s.shapes.add_picture(str(R / "dqn_learning_curve.png"), Inches(7.05), Inches(1.35),
                     width=Inches(5.9))
text(s, Inches(7.05), Inches(6.9), Inches(5.9), Inches(0.4),
     "Dueling DQN learning curve — stable convergence, no divergence",
     size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ================================================================ #
# 7 · RIGOR + REPRODUCE
# ================================================================ #
s = header(slide(), "Publication Rigor & How to Reproduce",
           "Statistically validated; one command regenerates everything")
bullets(s, Inches(0.55), Inches(1.35), Inches(6.3), Inches(5.5), [
    "5-seed evaluation | mean ± std across seeds 0–4 — DQN tops the table with tight variance (cost 29,141 ± 3,097)",
    "Wilcoxon signed-rank | significantly better than FCFS/RR/Greedy on all metrics (p < 0.0002, n = 75 pairs)",
    "4-way ablation | Dueling (+34% if removed) and mask (+37%) are the load-bearing ideas; PER ≈ neutral at this scale — reported honestly",
    "Dual datasets | Google Borg (cloud, light) + KTH SP2 (HPC, heavy) — robustness where heuristic rankings invert",
], size=12.5)
box = rect(s, Inches(7.15), Inches(1.5), Inches(5.6), Inches(4.6), RGBColor(0x0A, 0x1E, 0x33))
text(s, Inches(7.45), Inches(1.75), Inches(5.0), Inches(4.2),
     [("> reproduce everything", {"size": 13, "bold": True, "color": TEAL}),
      ("cd cloud_rl_xai_project_v2", {}),
      ("python run_experiment.py", {"color": RGBColor(0x7E, 0xE8, 0xB0)}),
      ("   # train + evaluate + explain  (~14 min)", {"size": 10.5, "color": MUTED}),
      ("python plots.py", {"color": RGBColor(0x7E, 0xE8, 0xB0)}),
      ("python build_dashboard.py", {"color": RGBColor(0x7E, 0xE8, 0xB0)}),
      ("python rigor.py --job aggregate", {"color": RGBColor(0x7E, 0xE8, 0xB0)}),
      ("", {}),
      ("Outputs:  dashboard.html · Figure 1 (PDF/SVG)", {"size": 11, "color": RGBColor(0xC9, 0xD6, 0xE8)}),
      ("decision_log.json (60 explained decisions)", {"size": 11, "color": RGBColor(0xC9, 0xD6, 0xE8)}),
      ("significance.json · ablation.json · manifest", {"size": 11, "color": RGBColor(0xC9, 0xD6, 0xE8)})],
     size=12, color=RGBColor(0xE8, 0xEE, 0xF5), leading=1.25)
text(s, Inches(7.15), Inches(6.3), Inches(5.6), Inches(0.8),
     "Live demo: dashboard Decision-Log tab — pick any row: the AI chose VM1 because vm1_util pushed its Q-value up (+2.59), verified faithful.",
     size=11.5, color=NAVY)

prs.save(OUT)
print(f"saved -> {OUT}")
